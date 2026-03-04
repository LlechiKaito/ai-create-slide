import base64
import io

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from backend.src.constants.slide import (
    ACCENT_BAR_HEIGHT_EMU,
    DEFAULT_CONTENT_GAP,
    DEFAULT_FONT_FAMILY,
    DEFAULT_FONT_SIZE_LEVEL,
    DEFAULT_IMAGE_SIZE,
    DIAGRAM_TYPES,
    PPTX_CHART_HEIGHT_INCHES,
    PPTX_CHART_LEFT_INCHES,
    PPTX_CHART_TOP_INCHES,
    PPTX_CHART_WIDTH_INCHES,
    PPTX_CONTENT_GAPS,
    PPTX_FONT_MAP,
    PPTX_FONT_SIZE_ACCENT,
    PPTX_FONT_SIZE_BODY,
    PPTX_FONT_SIZE_TITLE,
    PPTX_IMAGE_RIGHT_MARGIN,
    PPTX_IMAGE_SIZES,
    PPTX_TEXT_WIDTH_DEFAULT_INCHES,
    PPTX_TEXT_WIDTH_WITH_CHART_INCHES,
    SLIDE_HEIGHT_INCHES,
    SLIDE_WIDTH_INCHES,
    SUPPORTED_CHART_TYPES,
)
from backend.src.domain.commons.result import Result, success
from backend.src.domain.entities.slide.slide import Slide
from backend.src.domain.entities.slide.slide_deck import SlideDeck
from backend.src.domain.repositories.slide.slide_repository import SlideRepository
from backend.src.infrastructure.repositories.slide.pptx_diagram_renderer import (
    render_diagram,
)

DEFAULT_ACCENT = RGBColor(240, 130, 40)
DEFAULT_TEXT = RGBColor(50, 50, 50)
DEFAULT_BG = RGBColor(255, 255, 255)
GRAY = RGBColor(120, 120, 120)

BODY_LEFT_INCHES = 1.0
BODY_TOP_INCHES = 2.0
BULLET_SPACING_PT = 8

CHART_TYPE_MAP = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
}


def _hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _resolve_colors(color_config: dict | None) -> tuple[RGBColor, RGBColor, RGBColor]:
    if not color_config:
        return DEFAULT_ACCENT, DEFAULT_TEXT, DEFAULT_BG
    accent = _hex_to_rgb(color_config.get("accent", "#F08228"))
    text = _hex_to_rgb(color_config.get("text", "#323232"))
    bg = _hex_to_rgb(color_config.get("background", "#FFFFFF"))
    return accent, text, bg


def _resolve_font_sizes(cfg: dict) -> dict[str, int]:
    t = cfg.get("font_size_title", DEFAULT_FONT_SIZE_LEVEL)
    b = cfg.get("font_size_body", DEFAULT_FONT_SIZE_LEVEL)
    a = cfg.get("font_size_accent", DEFAULT_FONT_SIZE_LEVEL)
    title = PPTX_FONT_SIZE_TITLE.get(t, PPTX_FONT_SIZE_TITLE["medium"])
    body = PPTX_FONT_SIZE_BODY.get(b, PPTX_FONT_SIZE_BODY["medium"])
    accent = PPTX_FONT_SIZE_ACCENT.get(a, PPTX_FONT_SIZE_ACCENT["medium"])
    return {
        "main_title": int(title * 1.375),
        "title": title,
        "subtitle": accent,
        "body": body,
        "bullet": max(body - 2, 14),
        "bullet_marker": max(body - 6, 10),
    }


class PptxSlideRepository(SlideRepository):
    def generate_pptx(
        self, slide_deck: SlideDeck, color_config: dict | None = None
    ) -> Result[bytes, Exception]:
        accent, text_color, bg_color = _resolve_colors(color_config)
        cfg = color_config or {}
        font_name = PPTX_FONT_MAP.get(cfg.get("font_family", DEFAULT_FONT_FAMILY), "")
        image_size = cfg.get("image_size", DEFAULT_IMAGE_SIZE)
        content_gap = cfg.get("content_gap", DEFAULT_CONTENT_GAP)
        font_sizes = _resolve_font_sizes(cfg)

        prs = Presentation()
        prs.slide_width = Inches(SLIDE_WIDTH_INCHES)
        prs.slide_height = Inches(SLIDE_HEIGHT_INCHES)

        self._add_title_slide(
            prs, slide_deck, accent, text_color, bg_color, font_name, font_sizes,
        )

        for slide_entity in slide_deck.slides:
            self._add_content_slide(
                prs, slide_entity, accent, text_color, bg_color,
                font_name, image_size, content_gap, font_sizes,
            )

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return success(buffer.getvalue())

    def _set_slide_bg(self, slide: object, color: RGBColor) -> None:
        bg = slide.background  # type: ignore[attr-defined]
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_accent_bar(
        self, slide: object, accent: RGBColor,
        top: bool = False, bottom: bool = False,
    ) -> None:
        bar_h = Emu(ACCENT_BAR_HEIGHT_EMU)
        slide_w = Inches(SLIDE_WIDTH_INCHES)
        slide_h = Inches(SLIDE_HEIGHT_INCHES)
        if top:
            shape = slide.shapes.add_shape(1, 0, 0, slide_w, bar_h)  # type: ignore[attr-defined]
            shape.fill.solid()
            shape.fill.fore_color.rgb = accent
            shape.line.fill.background()
        if bottom:
            shape = slide.shapes.add_shape(1, 0, slide_h - bar_h, slide_w, bar_h)  # type: ignore[attr-defined]
            shape.fill.solid()
            shape.fill.fore_color.rgb = accent
            shape.line.fill.background()

    def _add_text_box(
        self,
        slide: object,
        left: Emu,
        top: Emu,
        width: Emu,
        height: Emu,
        text: str,
        font_size: int = 22,
        bold: bool = False,
        color: RGBColor = DEFAULT_TEXT,
        alignment: PP_ALIGN = PP_ALIGN.LEFT,
        font_name: str = "",
    ) -> None:
        txbox = slide.shapes.add_textbox(left, top, width, height)  # type: ignore[attr-defined]
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = alignment
        if font_name:
            p.font.name = font_name

    def _add_accent_line(
        self, slide: object, left: Emu, top: Emu, width: Emu, accent: RGBColor,
    ) -> None:
        shape = slide.shapes.add_shape(1, left, top, width, Emu(30000))  # type: ignore[attr-defined]
        shape.fill.solid()
        shape.fill.fore_color.rgb = accent
        shape.line.fill.background()

    def _add_title_slide(
        self, prs: Presentation, slide_deck: SlideDeck,
        accent: RGBColor, text_color: RGBColor, bg_color: RGBColor,
        font_name: str = "", font_sizes: dict | None = None,
    ) -> None:
        fs = font_sizes or _resolve_font_sizes({})
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_bg(slide, bg_color)
        self._add_accent_bar(slide, accent, top=True, bottom=True)

        slide_w = Inches(SLIDE_WIDTH_INCHES)

        self._add_text_box(
            slide, Inches(0), Inches(2.0), slide_w, Inches(1.0),
            slide_deck.title.value, fs["main_title"], True, text_color,
            PP_ALIGN.CENTER, font_name,
        )

        self._add_accent_line(slide, Inches(4.5), Inches(3.2), Inches(4.3), accent)

        if slide_deck.author:
            self._add_text_box(
                slide, Inches(0), Inches(3.5), slide_w, Inches(0.5),
                slide_deck.author, fs["subtitle"], False, GRAY,
                PP_ALIGN.CENTER, font_name,
            )

    def _add_image_to_slide(
        self, slide: object, image_data: str, image_size: str = DEFAULT_IMAGE_SIZE,
    ) -> None:
        size_cfg = PPTX_IMAGE_SIZES.get(image_size, PPTX_IMAGE_SIZES[DEFAULT_IMAGE_SIZE])
        img_w = size_cfg["width"]
        img_h = size_cfg["height"]
        img_left = SLIDE_WIDTH_INCHES - img_w - PPTX_IMAGE_RIGHT_MARGIN
        img_top = max(1.5, (SLIDE_HEIGHT_INCHES - img_h) / 2)

        raw = base64.b64decode(image_data)
        image_stream = io.BytesIO(raw)
        slide.shapes.add_picture(  # type: ignore[attr-defined]
            image_stream,
            Inches(img_left),
            Inches(img_top),
            Inches(img_w),
            Inches(img_h),
        )

    def _add_slide_header(
        self, slide: object, slide_entity: Slide,
        accent: RGBColor, text_color: RGBColor, font_name: str = "",
        font_sizes: dict | None = None,
    ) -> None:
        fs = font_sizes or _resolve_font_sizes({})
        slide_w = Inches(SLIDE_WIDTH_INCHES)
        if slide_entity.subtitle:
            self._add_text_box(
                slide, Inches(0), Inches(0.3), slide_w, Inches(0.4),
                slide_entity.subtitle, fs["subtitle"], False, accent,
                PP_ALIGN.CENTER, font_name,
            )
        self._add_text_box(
            slide, Inches(0), Inches(0.7), slide_w, Inches(0.7),
            slide_entity.title.value, fs["title"], True, text_color,
            PP_ALIGN.CENTER, font_name,
        )
        self._add_accent_line(slide, Inches(4.5), Inches(1.5), Inches(4.3), accent)

    def _add_slide_body(
        self, slide: object, slide_entity: Slide, text_w: Emu,
        accent: RGBColor, text_color: RGBColor, font_name: str = "",
        font_sizes: dict | None = None,
    ) -> None:
        fs = font_sizes or _resolve_font_sizes({})
        body_top = BODY_TOP_INCHES

        if slide_entity.content.value:
            self._add_text_box(
                slide, Inches(BODY_LEFT_INCHES), Inches(body_top), text_w, Inches(0.8),
                slide_entity.content.value, fs["body"], False, text_color,
                font_name=font_name,
            )
            body_top += 1.0

        if slide_entity.has_bullet_points():
            self._add_bullet_points(
                slide, slide_entity, text_w, body_top, accent, text_color,
                font_name, fs,
            )

    def _add_bullet_points(
        self, slide: object, slide_entity: Slide, text_w: Emu, top: float,
        accent: RGBColor, text_color: RGBColor, font_name: str = "",
        font_sizes: dict | None = None,
    ) -> None:
        fs = font_sizes or _resolve_font_sizes({})
        bullet_height = len(slide_entity.bullet_points) * 0.6 + 0.5
        txbox = slide.shapes.add_textbox(  # type: ignore[attr-defined]
            Inches(BODY_LEFT_INCHES), Inches(top), text_w, Inches(bullet_height),
        )
        tf = txbox.text_frame
        tf.word_wrap = True

        for i, point in enumerate(slide_entity.bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.space_after = Pt(BULLET_SPACING_PT)

            marker = p.add_run()
            marker.text = "\u25cf  "
            marker.font.size = Pt(fs["bullet_marker"])
            marker.font.color.rgb = accent
            marker.font.bold = True
            if font_name:
                marker.font.name = font_name

            text_run = p.add_run()
            text_run.text = point
            text_run.font.size = Pt(fs["bullet"])
            text_run.font.color.rgb = text_color
            if font_name:
                text_run.font.name = font_name

    def _add_chart_to_slide(
        self, slide: object, chart_data: dict, accent: RGBColor,
    ) -> None:
        chart_type_str = chart_data.get("chart_type", "column")
        if chart_type_str not in SUPPORTED_CHART_TYPES:
            chart_type_str = "column"
        xl_chart_type = CHART_TYPE_MAP[chart_type_str]

        categories = chart_data.get("categories", [])
        series_list = chart_data.get("series", [])
        if not categories or not series_list:
            return

        data = CategoryChartData()
        data.categories = categories
        for s in series_list:
            data.add_series(s.get("name", ""), s.get("values", []))

        chart_shape = slide.shapes.add_chart(  # type: ignore[attr-defined]
            xl_chart_type,
            Inches(PPTX_CHART_LEFT_INCHES),
            Inches(PPTX_CHART_TOP_INCHES),
            Inches(PPTX_CHART_WIDTH_INCHES),
            Inches(PPTX_CHART_HEIGHT_INCHES),
            data,
        )

        chart = chart_shape.chart
        chart.has_legend = len(series_list) > 1
        chart_title = chart_data.get("title", "")
        if chart_title:
            chart.has_title = True
            chart.chart_title.text_frame.paragraphs[0].text = chart_title

        for i, s in enumerate(chart.series):
            r, g, b = accent
            s.format.fill.solid()
            if i == 0:
                s.format.fill.fore_color.rgb = RGBColor(r, g, b)
            else:
                factor = 0.6 + (i * 0.15)
                s.format.fill.fore_color.rgb = RGBColor(
                    min(int(r * factor), 255),
                    min(int(g * factor), 255),
                    min(int(b * factor), 255),
                )

    def _add_content_slide(
        self, prs: Presentation, slide_entity: Slide,
        accent: RGBColor, text_color: RGBColor, bg_color: RGBColor,
        font_name: str = "",
        image_size: str = DEFAULT_IMAGE_SIZE,
        content_gap: str = DEFAULT_CONTENT_GAP,
        font_sizes: dict | None = None,
    ) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_bg(slide, bg_color)
        self._add_accent_bar(slide, accent, top=True, bottom=True)

        is_diagram = (
            slide_entity.has_chart()
            and slide_entity.chart_data
            and slide_entity.chart_data.get("chart_type") in DIAGRAM_TYPES
        )
        has_chart = (
            not is_diagram and slide_entity.has_chart() and slide_entity.chart_data
        )
        has_image = (
            not is_diagram and not has_chart and bool(slide_entity.image_data)
        )

        if has_image:
            size_cfg = PPTX_IMAGE_SIZES.get(image_size, PPTX_IMAGE_SIZES[DEFAULT_IMAGE_SIZE])
            gap_val = PPTX_CONTENT_GAPS.get(content_gap, PPTX_CONTENT_GAPS[DEFAULT_CONTENT_GAP])
            img_left = SLIDE_WIDTH_INCHES - size_cfg["width"] - PPTX_IMAGE_RIGHT_MARGIN
            text_w = Inches(img_left - gap_val - BODY_LEFT_INCHES)
        elif has_chart:
            text_w = Inches(PPTX_TEXT_WIDTH_WITH_CHART_INCHES)
        else:
            text_w = Inches(PPTX_TEXT_WIDTH_DEFAULT_INCHES)

        self._add_slide_header(slide, slide_entity, accent, text_color, font_name, font_sizes)
        self._add_slide_body(slide, slide_entity, text_w, accent, text_color, font_name, font_sizes)

        if is_diagram and slide_entity.chart_data:
            render_diagram(slide, slide_entity.chart_data, accent, text_color, font_name)
        elif has_chart:
            self._add_chart_to_slide(slide, slide_entity.chart_data, accent)
        elif has_image:
            self._add_image_to_slide(slide, slide_entity.image_data, image_size)
