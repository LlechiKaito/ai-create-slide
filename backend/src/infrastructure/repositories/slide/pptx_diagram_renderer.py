"""PPTX diagram renderer for shape-based diagrams."""
import math

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from backend.src.constants.slide import DIAGRAM_TYPES

TIMELINE_Y = 4.5
TIMELINE_LEFT = 1.2
TIMELINE_RIGHT = 12.1
TIMELINE_MARKER_EMU = 200000
FONT_PERIOD = 14
FONT_LABEL = 11

PROCESS_Y = 3.2
PROCESS_H = 1.5
PROCESS_LEFT = 1.2
PROCESS_RIGHT = 12.1
PROCESS_ARROW_W = 0.3

CYCLE_CX = 6.667
CYCLE_CY = 4.5
CYCLE_R = 2.2
CYCLE_NODE = 1.4

PYR_TOP = 2.8
PYR_BOTTOM = 6.8
PYR_CX = 6.667
PYR_MAX_W = 8.0

COMP_TOP = 2.5
COMP_COL_W = 5.0
COMP_GAP = 0.5
COMP_LEFT_X = 1.4

FONT_DIAGRAM = 12
FONT_DIAGRAM_SM = 10
FONT_COMP_TITLE = 16

WHITE = RGBColor(255, 255, 255)


def _add_textbox(
    slide: object, left: Emu, top: Emu, width: Emu, height: Emu,
    text: str, font_size: int, bold: bool, color: RGBColor,
    font_name: str = "", align: PP_ALIGN = PP_ALIGN.CENTER,
) -> None:
    txbox = slide.shapes.add_textbox(left, top, width, height)  # type: ignore[attr-defined]
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    if font_name:
        p.font.name = font_name


def _shade(accent: RGBColor, i: int) -> RGBColor:
    r, g, b = accent
    f = max(1.0 - i * 0.12, 0.3)
    return RGBColor(max(int(r * f), 0), max(int(g * f), 0), max(int(b * f), 0))


def render_diagram(
    slide: object, chart_data: dict,
    accent: RGBColor, text_color: RGBColor, font_name: str = "",
) -> None:
    dtype = chart_data.get("chart_type", "")
    if dtype not in DIAGRAM_TYPES:
        return
    renderer = _RENDERERS.get(dtype)
    if renderer:
        renderer(slide, chart_data, accent, text_color, font_name)


def _render_timeline(
    slide: object, chart_data: dict,
    accent: RGBColor, text_color: RGBColor, font_name: str,
) -> None:
    items = chart_data.get("items", [])
    if not items:
        return

    left = Inches(TIMELINE_LEFT)
    right = Inches(TIMELINE_RIGHT)
    line_y = Inches(TIMELINE_Y)
    line_w = right - left
    m_emu = Emu(TIMELINE_MARKER_EMU)
    half = Emu(TIMELINE_MARKER_EMU // 2)

    line = slide.shapes.add_shape(1, left, line_y, line_w, Emu(25000))  # type: ignore[attr-defined]
    line.fill.solid()
    line.fill.fore_color.rgb = accent
    line.line.fill.background()

    n = len(items)
    for i, item in enumerate(items):
        cx = (
            left + line_w // 2
            if n == 1
            else left + Emu(int(line_w * i / (n - 1)))
        )

        dot = slide.shapes.add_shape(9, cx - half, line_y - half, m_emu, m_emu)  # type: ignore[attr-defined]
        dot.fill.solid()
        dot.fill.fore_color.rgb = accent
        dot.line.fill.background()

        period = item.get("period", "")
        if period:
            _add_textbox(
                slide, cx - Inches(1.0), line_y - Inches(0.7),
                Inches(2.0), Inches(0.4), period,
                FONT_PERIOD, True, accent, font_name,
            )

        label = item.get("label", "")
        if label:
            _add_textbox(
                slide, cx - Inches(1.0), line_y + Inches(0.3),
                Inches(2.0), Inches(0.6), label,
                FONT_LABEL, False, text_color, font_name,
            )


def _render_process(
    slide: object, chart_data: dict,
    accent: RGBColor, text_color: RGBColor, font_name: str,
) -> None:
    items = chart_data.get("items", [])
    if not items:
        return

    n = len(items)
    total = PROCESS_RIGHT - PROCESS_LEFT
    arrow_w = PROCESS_ARROW_W if n > 1 else 0
    box_w = (total - arrow_w * max(n - 1, 0)) / max(n, 1)

    y = Inches(PROCESS_Y)
    h = Inches(PROCESS_H)

    for i, item in enumerate(items):
        x = Inches(PROCESS_LEFT + i * (box_w + arrow_w))

        box = slide.shapes.add_shape(1, x, y, Inches(box_w), h)  # type: ignore[attr-defined]
        box.fill.solid()
        box.fill.fore_color.rgb = _shade(accent, i)
        box.line.fill.background()

        label = item.get("label", "")
        if label:
            _add_textbox(
                slide, x, y + Inches(0.2), Inches(box_w), Inches(0.5),
                label, FONT_DIAGRAM, True, WHITE, font_name,
            )

        desc = item.get("description", "")
        if desc:
            _add_textbox(
                slide, x + Inches(0.1), y + Inches(0.7),
                Inches(box_w - 0.2), Inches(0.6),
                desc, FONT_DIAGRAM_SM, False, WHITE, font_name,
            )

        if i < n - 1:
            ax = x + Inches(box_w)
            arrow = slide.shapes.add_shape(  # type: ignore[attr-defined]
                5, ax + Inches(0.05), y + Inches(0.45),
                Inches(0.2), Inches(0.5),
            )
            arrow.rotation = 90.0
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = accent
            arrow.line.fill.background()


def _render_cycle(
    slide: object, chart_data: dict,
    accent: RGBColor, text_color: RGBColor, font_name: str,
) -> None:
    items = chart_data.get("items", [])
    if not items:
        return

    n = len(items)
    cx = Inches(CYCLE_CX)
    cy = Inches(CYCLE_CY)
    radius = Inches(CYCLE_R)
    node_s = Inches(CYCLE_NODE)
    half_n = Inches(CYCLE_NODE / 2)

    for i, item in enumerate(items):
        angle = 2 * math.pi * i / n - math.pi / 2
        nx = cx + Emu(int(radius * math.cos(angle))) - half_n
        ny = cy + Emu(int(radius * math.sin(angle))) - half_n

        circle = slide.shapes.add_shape(9, nx, ny, node_s, node_s)  # type: ignore[attr-defined]
        circle.fill.solid()
        circle.fill.fore_color.rgb = _shade(accent, i)
        circle.line.fill.background()

        label = item.get("label", "")
        if label:
            _add_textbox(
                slide, nx + Inches(0.1), ny + Inches(0.3),
                Inches(CYCLE_NODE - 0.2), Inches(0.8),
                label, FONT_LABEL, True, WHITE, font_name,
            )


def _render_pyramid(
    slide: object, chart_data: dict,
    accent: RGBColor, text_color: RGBColor, font_name: str,
) -> None:
    items = chart_data.get("items", [])
    if not items:
        return

    n = len(items)
    total_h = PYR_BOTTOM - PYR_TOP
    row_h = total_h / n

    for i, item in enumerate(items):
        ratio = (i + 1) / n
        w = PYR_MAX_W * (0.3 + 0.7 * ratio)
        x = Inches(PYR_CX - w / 2)
        y = Inches(PYR_TOP + i * row_h)
        h = Inches(row_h - 0.08)

        box = slide.shapes.add_shape(1, x, y, Inches(w), h)  # type: ignore[attr-defined]
        box.fill.solid()
        box.fill.fore_color.rgb = _shade(accent, i)
        box.line.fill.background()

        label = item.get("label", "")
        if label:
            _add_textbox(
                slide, x, y + Inches(0.1), Inches(w), Inches(row_h - 0.2),
                label, FONT_DIAGRAM, True, WHITE, font_name,
            )


def _render_funnel(
    slide: object, chart_data: dict,
    accent: RGBColor, text_color: RGBColor, font_name: str,
) -> None:
    items = chart_data.get("items", [])
    if not items:
        return

    n = len(items)
    total_h = PYR_BOTTOM - PYR_TOP
    row_h = total_h / n

    for i, item in enumerate(items):
        ratio = 1.0 - (i / max(n, 1))
        w = PYR_MAX_W * (0.3 + 0.7 * ratio)
        x = Inches(PYR_CX - w / 2)
        y = Inches(PYR_TOP + i * row_h)
        h = Inches(row_h - 0.08)

        box = slide.shapes.add_shape(1, x, y, Inches(w), h)  # type: ignore[attr-defined]
        box.fill.solid()
        box.fill.fore_color.rgb = _shade(accent, i)
        box.line.fill.background()

        label = item.get("label", "")
        if label:
            _add_textbox(
                slide, x, y + Inches(0.1), Inches(w), Inches(row_h - 0.2),
                label, FONT_DIAGRAM, True, WHITE, font_name,
            )


def _render_comparison(
    slide: object, chart_data: dict,
    accent: RGBColor, text_color: RGBColor, font_name: str,
) -> None:
    left_d = chart_data.get("left", {})
    right_d = chart_data.get("right", {})
    col_w = Inches(COMP_COL_W)
    lx = Inches(COMP_LEFT_X)
    rx = lx + col_w + Inches(COMP_GAP)
    ty = Inches(COMP_TOP)

    for data, x, shade_i in [(left_d, lx, 0), (right_d, rx, 2)]:
        title = data.get("title", "")
        if title:
            hdr = slide.shapes.add_shape(1, x, ty, col_w, Inches(0.8))  # type: ignore[attr-defined]
            hdr.fill.solid()
            hdr.fill.fore_color.rgb = _shade(accent, shade_i)
            hdr.line.fill.background()
            _add_textbox(
                slide, x, ty + Inches(0.1), col_w, Inches(0.6),
                title, FONT_COMP_TITLE, True, WHITE, font_name,
            )

        iy = ty + Inches(1.0)
        for item_text in data.get("items", []):
            _add_textbox(
                slide, x + Inches(0.3), iy,
                Inches(COMP_COL_W - 0.3), Inches(0.5),
                f"\u25cf {item_text}", FONT_DIAGRAM, False, text_color,
                font_name, PP_ALIGN.LEFT,
            )
            iy += Inches(0.6)


_RENDERERS = {
    "timeline": _render_timeline,
    "process": _render_process,
    "cycle": _render_cycle,
    "pyramid": _render_pyramid,
    "funnel": _render_funnel,
    "comparison": _render_comparison,
}
