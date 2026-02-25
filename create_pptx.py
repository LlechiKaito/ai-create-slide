from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

OUTPUT_DIR = Path("output/slides")
ASSETS_DIR = Path("output/assets")

# Colors
ORANGE = RGBColor(240, 130, 40)
DARK = RGBColor(50, 50, 50)
GRAY = RGBColor(120, 120, 120)
WHITE = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(245, 245, 245)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, top=False, bottom=False):
    """Add orange accent bars."""
    bar_h = Emu(80000)  # ~2mm
    if top:
        shape = slide.shapes.add_shape(1, 0, 0, SLIDE_W, bar_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = ORANGE
        shape.line.fill.background()
    if bottom:
        shape = slide.shapes.add_shape(1, 0, SLIDE_H - bar_h, SLIDE_W, bar_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = ORANGE
        shape.line.fill.background()


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name=""):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    if font_name:
        p.font.name = font_name
    return tf


def add_multi_text(slide, left, top, width, height, lines):
    """Add textbox with multiple styled lines.
    lines: list of (text, font_size, bold, color)
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, font_size, bold, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return tf


def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(1, left, top, width, Emu(30000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()


def add_card(slide, left, top, width, height, title, body_lines, title_color=ORANGE):
    """Add a rounded card with title and body."""
    # Card background
    shape = slide.shapes.add_shape(
        5,  # rounded rectangle
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.fill.background()
    shape.shadow.inherit = False

    # Title
    add_text_box(slide, left + Inches(0.3), top + Inches(0.2),
                 width - Inches(0.6), Inches(0.5),
                 title, font_size=20, bold=True, color=title_color)

    # Body
    y = top + Inches(0.8)
    for line in body_lines:
        add_text_box(slide, left + Inches(0.3), y,
                     width - Inches(0.6), Inches(0.4),
                     line, font_size=14, color=DARK)
        y += Inches(0.35)


# ── Slides ──

def slide_1_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)
    add_accent_bar(slide, top=True, bottom=True)

    # Subtitle
    add_text_box(slide, Inches(0), Inches(1.8), SLIDE_W, Inches(0.5),
                 "SELF INTRODUCTION", font_size=20, bold=False,
                 color=ORANGE, alignment=PP_ALIGN.CENTER)

    # Accent line
    add_accent_line(slide, Inches(4.5), Inches(2.4), Inches(4.3))

    # Main title
    add_text_box(slide, Inches(0), Inches(2.6), SLIDE_W, Inches(1),
                 "自己紹介", font_size=48, bold=True,
                 color=DARK, alignment=PP_ALIGN.CENTER)

    # Name
    add_text_box(slide, Inches(0), Inches(3.7), SLIDE_W, Inches(0.7),
                 "一条海音", font_size=36, bold=False,
                 color=DARK, alignment=PP_ALIGN.CENTER)

    # Accent line
    add_accent_line(slide, Inches(4.5), Inches(4.5), Inches(4.3))

    # Catchphrase
    add_text_box(slide, Inches(0), Inches(4.7), SLIDE_W, Inches(0.6),
                 "webエンジニア × AI で未来をつくる", font_size=24, bold=False,
                 color=ORANGE, alignment=PP_ALIGN.CENTER)


def slide_2_skills(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, top=True, bottom=True)

    # Title block
    add_text_box(slide, Inches(0), Inches(0.4), SLIDE_W, Inches(0.4),
                 "CAREER & SKILLS", font_size=18, color=ORANGE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0), Inches(0.8), SLIDE_W, Inches(0.7),
                 "経歴・スキル", font_size=36, bold=True,
                 color=DARK, alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, Inches(4.5), Inches(1.5), Inches(4.3))

    # Left card - Career
    add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.2),
             "経歴", [
                 "芝浦工業大学 システム理工学部",
                 "電子情報システム学科 卒業",
                 "AIによる画像認識の研究に従事",
                 "",
                 "株式会社Betterboundにて9ヶ月間",
                 "長期インターン（yarikiri開発）",
                 "",
                 "現在は大企業でtoCサービスの自社開発",
                 "要件定義〜実装〜デプロイ〜テストまで",
                 "横断的に担当",
                 "AI推進として業務改善にも取り組む",
             ])

    # Right card - Tech Stack (separate label/value for alignment)
    card_left = Inches(7.0)
    card_top = Inches(1.8)
    card_w = Inches(5.5)
    card_h = Inches(5.2)

    shape = slide.shapes.add_shape(5, card_left, card_top, card_w, card_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.fill.background()
    shape.shadow.inherit = False

    add_text_box(slide, card_left + Inches(0.3), card_top + Inches(0.2),
                 card_w - Inches(0.6), Inches(0.5),
                 "技術スタック", font_size=20, bold=True, color=ORANGE)

    label_x = card_left + Inches(0.2)
    label_w = Inches(2.2)
    value_x = card_left + Inches(2.5)
    value_w = card_w - Inches(2.7)
    row_h = Inches(0.35)
    y = card_top + Inches(0.8)

    tech_items = [
        ("フロントエンド", "Vue.js / Nuxt.js /"),
        (None,            "React / TypeScript / Next.js"),
        ("バックエンド",   "C#(.NET) /"),
        (None,            "TypeScript(Express.js)"),
        ("インフラ",       "AWS"),
        ("AI",             "Claude CodeでのAI駆動開発"),
        ("その他",         "Git / Linux / Docker"),
    ]

    for label, value in tech_items:
        if label:
            add_text_box(slide, label_x, y, label_w, Inches(0.4),
                         label + " :", font_size=14, color=DARK,
                         alignment=PP_ALIGN.RIGHT)
        add_text_box(slide, value_x, y, value_w, Inches(0.4),
                     value, font_size=14, color=DARK)
        y += row_h

    # Image
    img_path = ASSETS_DIR / "work_miku.png"
    if img_path.exists():
        slide.shapes.add_picture(
            str(img_path), Inches(10.5), Inches(5.2), Inches(2.2), Inches(2.2))


def slide_3_work(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, top=True, bottom=True)

    # Title block
    add_text_box(slide, Inches(0), Inches(0.4), SLIDE_W, Inches(0.4),
                 "CURRENT WORK", font_size=18, color=ORANGE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0), Inches(0.8), SLIDE_W, Inches(0.7),
                 "今やっていること", font_size=36, bold=True,
                 color=DARK, alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, Inches(4.5), Inches(1.5), Inches(4.3))

    # Main content
    add_text_box(slide, Inches(1.2), Inches(2.2), Inches(7), Inches(0.5),
                 "webエンジニアとして", font_size=28, bold=True, color=ORANGE)

    lines = [
        ("・大規模Webアプリケーションの設計・開発", 18, False, DARK),
        ("・フロントエンドからバックエンドまで幅広く対応", 18, False, DARK),
        ("・チームメンバーと協力しながらプロダクトを推進", 18, False, DARK),
        ("", 10, False, DARK),
        ("AI活用にも積極的に取り組んでいます", 22, True, ORANGE),
        ("・Claude Codeを活用したAI駆動開発", 18, False, DARK),
        ("・AIを使った開発ワークフローの改善", 18, False, DARK),
    ]
    add_multi_text(slide, Inches(1.2), Inches(2.9), Inches(7), Inches(4), lines)

    # Banner image (isekai anime style)
    img_path = ASSETS_DIR / "isekai_work.png"
    if img_path.exists():
        slide.shapes.add_picture(
            str(img_path), Inches(9.0), Inches(2.2), Inches(3.5), Inches(2.0))


def slide_4_anime(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, top=True, bottom=True)

    # Title block
    add_text_box(slide, Inches(0), Inches(0.4), SLIDE_W, Inches(0.4),
                 "HOBBY 1", font_size=18, color=ORANGE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0), Inches(0.8), SLIDE_W, Inches(0.7),
                 "趣味その１ ─ アニメ", font_size=36, bold=True,
                 color=DARK, alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, Inches(4.5), Inches(1.5), Inches(4.3))

    # Image (16:9 ratio preserved)
    img_w = Inches(6.5)
    img_h = Inches(3.656)  # 6.5 * 9/16
    img_path = ASSETS_DIR / "anime3.png"
    if img_path.exists():
        slide.shapes.add_picture(
            str(img_path), Inches(0.8), Inches(1.8), img_w, img_h)

    # Text on the right
    tx = Inches(7.8)
    tw = Inches(4.8)

    add_text_box(slide, tx, Inches(1.8), tw, Inches(0.5),
                 "深夜アニメが大好き", font_size=24, bold=True, color=ORANGE)

    add_text_box(slide, tx, Inches(2.4), tw, Inches(0.4),
                 "特に転生系アニメにハマっています！", font_size=16, color=DARK)

    add_text_box(slide, tx, Inches(3.2), tw, Inches(0.5),
                 "お気に入り作品", font_size=24, bold=True, color=ORANGE)

    favorites = [
        ("・転生したら剣でした", 16, False, DARK),
        ("・転生したらスライムだった件", 16, False, DARK),
        ("・無職転生", 16, False, DARK),
        ("・戦姫絶唱シンフォギア", 16, False, DARK),
    ]
    add_multi_text(slide, tx, Inches(3.8), tw, Inches(2.5), favorites)


def slide_5_vocaloid(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, top=True, bottom=True)

    # Title block
    add_text_box(slide, Inches(0), Inches(0.4), SLIDE_W, Inches(0.4),
                 "HOBBY 2", font_size=18, color=ORANGE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0), Inches(0.8), SLIDE_W, Inches(0.7),
                 "趣味その２ ─ ボカロ", font_size=36, bold=True,
                 color=DARK, alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, Inches(4.5), Inches(1.5), Inches(4.3))

    # Image (16:9 ratio preserved)
    img_w = Inches(6.5)
    img_h = Inches(3.656)  # 6.5 * 9/16
    img_path = ASSETS_DIR / "vocaloid3.png"
    if img_path.exists():
        slide.shapes.add_picture(
            str(img_path), Inches(0.8), Inches(1.8), img_w, img_h)

    # Text on the right
    tx = Inches(7.8)
    tw = Inches(4.8)

    add_text_box(slide, tx, Inches(1.8), tw, Inches(0.5),
                 "初音ミクが大好き！", font_size=24, bold=True, color=ORANGE)

    add_text_box(slide, tx, Inches(2.4), tw, Inches(0.5),
                 "毎年ライブに参加しています。", font_size=16, color=DARK)

    add_text_box(slide, tx, Inches(3.2), tw, Inches(0.5),
                 "参戦イベント", font_size=24, bold=True, color=ORANGE)

    events = [
        ("・マジカルミライ（マジミラ）", 16, False, DARK),
        ("・セカライ", 16, False, DARK),
    ]
    add_multi_text(slide, tx, Inches(3.8), tw, Inches(1.5), events)


def slide_6_workstyle(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, top=True, bottom=True)

    # Title block
    add_text_box(slide, Inches(0), Inches(0.4), SLIDE_W, Inches(0.4),
                 "WORK STYLE", font_size=18, color=ORANGE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0), Inches(0.8), SLIDE_W, Inches(0.7),
                 "一緒に仕事するときの特徴", font_size=36, bold=True,
                 color=DARK, alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, Inches(4.5), Inches(1.5), Inches(4.3))

    # Cards in 3 columns
    card_w = Inches(3.5)
    card_h = Inches(4.8)
    gap = Inches(0.5)
    start_x = Inches(0.9)
    card_y = Inches(2.0)

    add_card(slide, start_x, card_y, card_w, card_h,
             "開発スタイル", [
                 "設計から実装まで一気通貫で",
                 "取り組むのが好きです。",
                 "",
                 "コードレビューは丁寧に",
                 "建設的なフィードバックを",
                 "心がけています。",
             ])

    add_card(slide, start_x + card_w + gap, card_y, card_w, card_h,
             "コミュニケーション", [
                 "テキストでの情報共有を",
                 "大事にしています。",
                 "",
                 "困ったときは早めに相談！",
                 "気軽に声をかけてください。",
             ])

    add_card(slide, start_x + (card_w + gap) * 2, card_y, card_w, card_h,
             "得意なこと", [
                 "新しい技術のキャッチアップ",
                 "",
                 "Claude Codeを使った",
                 "AI駆動開発",
                 "",
                 "フロントエンド",
                 "バックエンド",
                 "インフラ横断的な開発",
             ])


def slide_7_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, top=True, bottom=True)

    # Title block
    add_text_box(slide, Inches(0), Inches(0.4), SLIDE_W, Inches(0.4),
                 "SUMMARY", font_size=18, color=ORANGE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0), Inches(0.8), SLIDE_W, Inches(0.7),
                 "まとめ", font_size=36, bold=True,
                 color=DARK, alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, Inches(4.5), Inches(1.5), Inches(4.3))

    # Summary items
    items = [
        ("webエンジニア", "大企業でフルスタックに活躍中"),
        ("AI活用", "Claude CodeでAI駆動開発に取り組んでいます"),
        ("アニメ", "転生系アニメ・シンフォギアが好き"),
        ("ボカロ", "初音ミク推し / マジミラ・セカライ参戦"),
    ]

    y = Inches(2.0)
    for label, desc in items:
        add_text_box(slide, Inches(2.5), y, Inches(3), Inches(0.5),
                     label, font_size=22, bold=True, color=ORANGE,
                     alignment=PP_ALIGN.RIGHT)
        add_text_box(slide, Inches(5.8), y + Inches(0.03), Inches(6), Inches(0.5),
                     desc, font_size=18, color=DARK)
        y += Inches(0.65)

    # Closing message
    y += Inches(0.5)
    add_text_box(slide, Inches(0), y, SLIDE_W, Inches(0.7),
                 "気軽に話しかけてください！", font_size=28, bold=True,
                 color=ORANGE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0), y + Inches(0.7), SLIDE_W, Inches(0.5),
                 "Thank you!", font_size=32, bold=True,
                 color=DARK, alignment=PP_ALIGN.CENTER)


def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_1_title(prs)
    slide_2_skills(prs)
    slide_3_work(prs)
    slide_4_anime(prs)
    slide_5_vocaloid(prs)
    slide_6_workstyle(prs)
    slide_7_summary(prs)

    output_path = OUTPUT_DIR / "self_introduction.pptx"
    prs.save(str(output_path))
    print(f"Created: {output_path}")
    print(f"Done! 7 slides saved.")


if __name__ == "__main__":
    main()
